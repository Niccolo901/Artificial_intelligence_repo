from pptx import Presentation
from pptx.util import Inches
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
data = pd.read_csv("housing.csv")

# Initialize the presentation
prs = Presentation()

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Slide 1: Title Slide with Customizations
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Add a background image (optional)
background_img = "plots/losa_angeles_lanscape.png"  # Replace with your background image path
slide.shapes.add_picture(background_img, 0, 0, width=prs.slide_width, height=prs.slide_height)

# Add a title text box
title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True

# Customize title font and content
title = title_frame.add_paragraph()
title.text = "Housing Price Analysis: Clustering and Prediction"
title.font.bold = True
title.font.size = Pt(44)
title.font.color.rgb = RGBColor(255, 255, 255)  # White text
title.alignment = 1  # Center alignment

# Add a subtitle text box
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True

# Customize subtitle font and content
subtitle = subtitle_frame.add_paragraph()
subtitle.text = "A Machine Learning Approach to Real Estate Insights"
subtitle.font.size = Pt(28)
subtitle.font.color.rgb = RGBColor(200, 200, 200)  # Light gray text
subtitle.alignment = 1  # Center alignment


# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Introduction"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT  # Corrected alignment for the title

# Problem Statement
problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2))
problem_frame = problem_box.text_frame
problem_frame.word_wrap = True

problem_title = problem_frame.add_paragraph()
problem_title.text = "Problem Statement:"
problem_title.font.bold = True
problem_title.font.size = Pt(24)
problem_title.font.color.rgb = RGBColor(0, 102, 204)  # Lighter blue

problem_content = problem_frame.add_paragraph()
problem_content.text = (
    "- What are the factors affecting house prices?\n"
    "- How can clustering and prediction help in real estate analysis?"
)
problem_content.font.size = Pt(18)
problem_content.font.color.rgb = RGBColor(0, 0, 0)  # Black

# Objectives
objective_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(8), Inches(2))
objective_frame = objective_box.text_frame
objective_frame.word_wrap = True

objective_title = objective_frame.add_paragraph()
objective_title.text = "Objectives:"
objective_title.font.bold = True
objective_title.font.size = Pt(24)
objective_title.font.color.rgb = RGBColor(0, 102, 204)  # Lighter blue

objective_content = objective_frame.add_paragraph()
objective_content.text = (
    "- Explore the dataset and identify trends.\n"
    "- Perform clustering to uncover patterns.\n"
    "- Build models to predict house prices."
)
objective_content.font.size = Pt(18)
objective_content.font.color.rgb = RGBColor(0, 0, 0)  # Black


# Slide 3: Overview of the Dataset
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Overview of the Dataset"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Dataset Shape Section
shape_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1))
shape_frame = shape_box.text_frame
shape_frame.word_wrap = True

shape_title = shape_frame.add_paragraph()
shape_title.text = "Dataset Shape:"
shape_title.font.bold = True
shape_title.font.size = Pt(24)
shape_title.font.color.rgb = RGBColor(0, 102, 204)

shape_content = shape_frame.add_paragraph()
shape_content.text = "- Total Rows: 20,640\n- Total Columns: 10"
shape_content.font.size = Pt(18)
shape_content.font.color.rgb = RGBColor(0, 0, 0)

# Key Features Section
features_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(8), Inches(2))
features_frame = features_box.text_frame
features_frame.word_wrap = True

features_title = features_frame.add_paragraph()
features_title.text = "Key Features:"
features_title.font.bold = True
features_title.font.size = Pt(24)
features_title.font.color.rgb = RGBColor(0, 102, 204)

features_content = features_frame.add_paragraph()
features_content.text = (
    "1. longitude: How far west a house is.\n"
    "2. latitude: How far north a house is.\n"
    "3. housingMedianAge: Median age of a house within a block.\n"
    "4. totalRooms: Total number of rooms within a block.\n"
    "5. totalBedrooms: Total number of bedrooms within a block.\n"
    "6. population: Total population within a block.\n"
    "7. households: Total households within a block.\n"
    "8. medianIncome: Median income for households.\n"
    "9. medianHouseValue: Median house value for households.\n"
    "10. oceanProximity: Proximity of the house to the ocean (categorical)."
)
features_content.font.size = Pt(18)
features_content.font.color.rgb = RGBColor(0, 0, 0)

# Dataset Head Section
head_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(8), Inches(1))
head_frame = head_box.text_frame
head_frame.word_wrap = True

head_title = head_frame.add_paragraph()
head_title.text = "Sample Data (First 5 Rows):"
head_title.font.bold = True
head_title.font.size = Pt(24)
head_title.font.color.rgb = RGBColor(0, 102, 204)

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Slide 4: Dataset Head Table
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Sample Data (First 5 Rows)"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Add Dataset Head Table
rows, cols = data.head().shape  # Get number of rows and columns
table = slide.shapes.add_table(rows=rows + 1, cols=cols, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(3)).table

# Set Table Headers
for col_idx, col_name in enumerate(data.columns):
    table.cell(0, col_idx).text = col_name  # Add column names
    table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True  # Bold header
    table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)  # Header font size
    table.cell(0, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    table.cell(0, col_idx).fill.solid()
    table.cell(0, col_idx).fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue background

# Populate Table Rows with Data
for row_idx, row_data in enumerate(data.head().itertuples(index=False), start=1):
    for col_idx, value in enumerate(row_data):
        table.cell(row_idx, col_idx).text = str(value)  # Add row values
        table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)  # Row font size
        table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Adjust Table Design
for row in range(rows + 1):
    for col in range(cols):
        cell = table.cell(row, col)
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Slide 5: Missing Values
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Missing Values in the Dataset"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "The dataset contains missing values in the 'total_bedrooms' column, "
    "with 207 missing entries out of 20,640 rows. Handling missing values is essential "
    "to ensure the quality and reliability of the analysis."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add Heatmap Image
heatmap_path = "plots/missing_values_heatmap.png"  # Path to the uploaded image
slide.shapes.add_picture(heatmap_path, Inches(0.5), Inches(3), width=Inches(9), height=Inches(4))


# Slide: Dataset Normalization and Encoding
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Dataset Normalization and Encoding"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "We normalized the dataset to bring all numerical features to the same scale and "
    "applied one-hot encoding to transform the categorical variable 'ocean_proximity' into "
    "numeric binary variables.\n\n"
    "The processed dataset is now ready for machine learning models and has been saved as 'processed_housing.csv'."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add Processed Dataset Table
# Add a small portion of the processed dataset as a table
processed_data_sample = [
    ["longitude", "latitude", "housing_median_age", "total_rooms", "total_bedrooms",
     "population", "households", "median_income", "ocean_proximity_INLAND", "ocean_proximity_NEAR BAY"],
    [-122.23, 37.88, 0.982163, -0.803813, -0.970325, -0.973320, -0.976833, 2.345163, False, True],
    [-122.22, 37.86, -0.606210, 2.042130, 1.348276, 0.861339, 1.670373, 2.332632, False, True],
    [-122.24, 37.85, 1.855769, -0.535189, -0.825561, -0.819769, -0.843427, 1.782939, False, True],
    [-122.25, 37.85, 1.855769, -0.623510, -0.718768, -0.765056, -0.733562, 0.932970, False, True],
    [-122.25, 37.85, 1.855769, -0.461970, -0.611974, -0.758879, -0.628930, -0.013143, False, True],
]

rows, cols = len(processed_data_sample), len(processed_data_sample[0])
table = slide.shapes.add_table(rows=rows, cols=cols, left=Inches(0.5), top=Inches(3.2), width=Inches(9), height=Inches(2.5)).table

# Fill Table Header
for col_idx, header in enumerate(processed_data_sample[0]):
    table.cell(0, col_idx).text = header
    table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True
    table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)
    table.cell(0, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    table.cell(0, col_idx).fill.solid()
    table.cell(0, col_idx).fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue background

# Fill Table Rows
for row_idx, row_data in enumerate(processed_data_sample[1:], start=1):
    for col_idx, value in enumerate(row_data):
        table.cell(row_idx, col_idx).text = str(value)
        table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)
        table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Slide: Map of Houses
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Map of Houses in the Dataset"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "The map below shows the geographic distribution of houses in California included in the dataset. "
    "Each point represents a house, with longitude and latitude used to plot their locations."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add Map Image
map_path = "plots/map_plot.png"  # Path to the uploaded map image
slide.shapes.add_picture(map_path, Inches(0.5), Inches(3), width=Inches(9), height=Inches(5))

# Slide: Correlation Matrix
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Correlation Matrix of Numeric Features"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "The correlation matrix provides an overview of relationships between numeric features in the dataset. "
    "It highlights significant correlations, such as the strong positive relationship between 'median_income' and 'median_house_value'."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add Correlation Matrix Image
matrix_path = "plots/correlation_matrix.png"  # Path to the uploaded image
slide.shapes.add_picture(matrix_path, Inches(0.5), Inches(3), width=Inches(9), height=Inches(5))

# Slide: Distributions and Boxplots
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Distributions and Boxplots of Key Features"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "The histograms and boxplots below provide insights into the distributions of key features:\n"
    "- `median_house_value`: Shows skewness and a cap at 500,000.\n"
    "- `median_income`: Indicates the range and presence of outliers."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add Images
# Add the first image (median_house_value distribution and boxplot)
house_value_path = "plots/media_house_values.png"  # Path to the first uploaded image
slide.shapes.add_picture(house_value_path, Inches(0.5), Inches(3), width=Inches(4.25), height=Inches(4))

# Add the second image (median_income distribution and boxplot)
income_path = "plots/median_income_boxplot.png"  # Path to the second uploaded image
slide.shapes.add_picture(income_path, Inches(4.75), Inches(3), width=Inches(4.25), height=Inches(4))


# Slide: Principal Component Analysis (PCA)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Principal Component Analysis (PCA)"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "Using PCA to reduce dimensionality:\n"
    "- Non-numeric columns converted where possible.\n"
    "- Kaiser Rule applied: Retaining 3 components (eigenvalues > 1).\n"
    "- Explained Variance Ratio (Selected Components): [0.4348, 0.2136, 0.1886]"
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add PCA Loadings Table
pca_loadings = [
    ["Feature", "PC1", "PC2", "PC3"],
    ["longitude", "0.0762", "-0.6620", "-0.2480"],
    ["latitude", "-0.0759", "0.6905", "0.1278"],
    ["housing_median_age", "-0.2164", "0.0009", "0.0496"],
    ["total_rooms", "0.4844", "0.0607", "0.0850"],
    ["total_bedrooms", "0.4895", "0.0721", "-0.0412"],
    ["population", "0.4701", "0.0461", "-0.0780"],
    ["households", "0.4910", "0.0718", "-0.0265"],
    ["median_income", "0.0558", "-0.1783", "0.6711"],
    ["median_house_value", "0.0452", "-0.1922", "0.6735"],
]

rows, cols = len(pca_loadings), len(pca_loadings[0])
table = slide.shapes.add_table(rows=rows, cols=cols, left=Inches(0.5), top=Inches(3.5), width=Inches(9), height=Inches(3.5)).table

# Set Table Headers
for col_idx, header in enumerate(pca_loadings[0]):
    table.cell(0, col_idx).text = header
    table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True
    table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(12)
    table.cell(0, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    table.cell(0, col_idx).fill.solid()
    table.cell(0, col_idx).fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue background

# Populate Table Rows with PCA Loadings
for row_idx, row_data in enumerate(pca_loadings[1:], start=1):
    for col_idx, value in enumerate(row_data):
        table.cell(row_idx, col_idx).text = value
        table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.size = Pt(10)
        table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Slide: Scree Plot
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Scree Plot (Eigenvalues)"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "The scree plot shows the eigenvalues of the principal components. "
    "Using the Kaiser Rule (eigenvalue > 1), we retain 3 components that capture "
    "the majority of the dataset's variance. This helps reduce dimensionality while "
    "maintaining important information."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add Scree Plot Image
scree_plot_path = "plots/scree_plot.png"  # Path to the uploaded scree plot image
slide.shapes.add_picture(scree_plot_path, Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(4))

# Slide: PCA Biplot
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "PCA Biplot"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "The PCA Biplot visualizes the dataset in the space of the first two principal components. "
    "It also shows the direction and magnitude of each original feature's contribution to these components. "
    "This helps to understand the influence of variables on the reduced dimensions."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add Biplot Image
biplot_path = "plots/biplot.png"  # Path to the uploaded biplot image
slide.shapes.add_picture(biplot_path, Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(4))

# Slide: Clustering Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Clustering Overview"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(3))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "Principal Component Analysis (PCA) was used to reduce dimensionality:\n"
    "- Explained Variance Ratio of Selected Components: [0.435, 0.214, 0.189]\n\n"
    "KMeans clustering was performed using PCA-reduced data:\n"
    "- The optimal number of clusters was determined by silhouette scores.\n"
    "- The silhouette scores for different cluster numbers are evaluated."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text


# Slide: Silhouette Scores
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Silhouette Scores and Optimal Clusters"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "The silhouette score measures the quality of clustering:\n"
    "- A higher silhouette score indicates better cluster separation.\n\n"
    "Results of silhouette analysis:\n"
    "Number of clusters: 2, Silhouette score: 0.312\n"
    "Number of clusters: 3, Silhouette score: 0.361\n"
    "Number of clusters: 4, Silhouette score: 0.361\n"
    "Number of clusters: 5, Silhouette score: 0.325\n"
    "Number of clusters: 6, Silhouette score: 0.340\n"
    "Number of clusters: 7, Silhouette score: 0.324\n"
    "Number of clusters: 8, Silhouette score: 0.296\n"
    "Number of clusters: 9, Silhouette score: 0.310\n"
    "Number of clusters: 10, Silhouette score: 0.307"
)
content.font.size = Pt(16)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Slide: Map of Houses with Clusters
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

# Title for the Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Map of Houses with Clusters"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Explanation Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "The map visualizes the clusters identified by KMeans in the context of their geographical locations.\n\n"
    "- Each house is color-coded according to its cluster.\n"
    "- Geographic patterns in clustering can provide insights into regional trends."
)
content.font.size = Pt(18)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add Map Clustering Plot
clustering_map_path = "plots/clustering_map_plo.png"  # Path to the map clustering image
slide.shapes.add_picture(clustering_map_path, Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(4))



# Add Silhouette Plot Image
silhouette_plot_path = "plots/silouhette_score_plot.png"  # Path to the silhouette score plot image
slide.shapes.add_picture(silhouette_plot_path, Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(4))


#Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Supervised Analysis"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Content
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "In the Supervised Analysis, we aim to predict housing prices based on the features in the dataset.\n\n"
    "Techniques used:\n"
    "• Linear Regression\n"
    "• Random Forest\n"
    "• Neural Network (TensorFlow)\n"
    "• XGBoost"
)
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add a blank slide for Linear Regression
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Linear Regression Results"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Content/Description
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content = content_frame.add_paragraph()
content.text = (
    "Linear regression was performed to predict house prices using the dataset's features.\n\n"
    "Key Results:\n"
    "• Coefficients provide insights into feature importance.\n"
    "• Positive coefficients indicate a direct relationship, while negative coefficients indicate an inverse relationship.\n"
    "• The plot shows the relative magnitude and direction of coefficients."
)
content.font.size = Pt(20)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add the image for coefficients
img_path = "plots/OLS_features.png"  # Replace with your actual file path
slide.shapes.add_picture(img_path, Inches(0.5), Inches(3.5), Inches(8), Inches(4))


# Add a new slide for Linear Regression Metrics
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Linear Regression: Performance Metrics"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Add metrics as text
metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
metrics_frame = metrics_box.text_frame
metrics_frame.word_wrap = True

metrics_content = metrics_frame.add_paragraph()
metrics_content.text = (
    "Linear Regression Model Performance:\n\n"
    "• Mean Absolute Error (MAE): 50,413.43\n"
    "• Mean Squared Error (MSE): 4,802,173,538.60\n"
    "• Root Mean Squared Error (RMSE): 69,297.72\n"
    "• R-squared (R²): 0.6488\n"
    "• Adjusted R-squared: 0.6478\n"
)
metrics_content.font.size = Pt(20)
metrics_content.font.color.rgb = RGBColor(0, 0, 0)  # Black text


# Add a new slide for Random Forest
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Random Forest: Performance Metrics"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Add metrics as text
metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
metrics_frame = metrics_box.text_frame
metrics_frame.word_wrap = True

metrics_content = metrics_frame.add_paragraph()
metrics_content.text = (
    "Random Forest Model Performance:\n\n"
    "• Best Parameters: {'max_depth': 20, 'min_samples_split': 2, 'n_estimators': 500}\n"
    "• Mean Absolute Error (MAE): 31,641.57\n"
    "• Mean Squared Error (MSE): 2,375,466,293.68\n"
    "• Root Mean Squared Error (RMSE): 48,738.76\n"
    "• R-squared (R²): 0.8263\n"
    "• Adjusted R-squared: 0.8258\n"
)
metrics_content.font.size = Pt(20)
metrics_content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Insert the feature importance image
img_path = "plots/rf_feature_importance.png"  # Update the file path if necessary
slide.shapes.add_picture(img_path, Inches(0.5), Inches(4.5), Inches(8), Inches(3))

# Add a new slide for Neural Network
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Neural Network: Performance Metrics"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Add metrics as text
metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(3))
metrics_frame = metrics_box.text_frame
metrics_frame.word_wrap = True

metrics_content = metrics_frame.add_paragraph()
metrics_content.text = (
    "Neural Network Model Performance:\n\n"
    "• Mean Absolute Error (MAE): 43,290.98\n"
    "• Mean Squared Error (MSE): 3,882,304,308.67\n"
    "• Root Mean Squared Error (RMSE): 62,308.14\n"
    "• R-squared (R²): 0.7161\n"
    "• Adjusted R-squared: 0.7153\n"
)
metrics_content.font.size = Pt(20)
metrics_content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Insert the Training and Validation Loss plot
img_path = "plots/Train_val_loss.png"  # Update the file path if necessary
slide.shapes.add_picture(img_path, Inches(0.5), Inches(4.5), Inches(8), Inches(3))


# Add a new slide for XGBoost
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "XGBoost: Performance Metrics and Feature Importance"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Add metrics as text
metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(3))
metrics_frame = metrics_box.text_frame
metrics_frame.word_wrap = True

metrics_content = metrics_frame.add_paragraph()
metrics_content.text = (
    "XGBoost Model Performance:\n\n"
    "• Mean Absolute Error (MAE): 30,660.61\n"
    "• Mean Squared Error (MSE): 2,204,422,834.38\n"
    "• Root Mean Squared Error (RMSE): 46,951.28\n"
    "• R-squared (R²): 0.8388\n"
    "• Adjusted R-squared: 0.8383\n\n"
    "Best Parameters:\n"
    "• colsample_bytree: 1.0\n"
    "• learning_rate: 0.1\n"
    "• max_depth: 5\n"
    "• n_estimators: 500\n"
    "• subsample: 0.8"
)
metrics_content.font.size = Pt(20)
metrics_content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Insert the Feature Importance plot
img_path = "plots/Xgboost_features.png"  # Update the file path if necessary
slide.shapes.add_picture(img_path, Inches(0.5), Inches(4.5), Inches(8), Inches(3))


# Add a new slide for Model Evaluation Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Model Evaluation Summary"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Add the Summary Table as text
summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(3))
summary_frame = summary_box.text_frame
summary_frame.word_wrap = True

summary_content = summary_frame.add_paragraph()
summary_content.text = (
    "Key Model Performance Metrics:\n\n"
    "1. Linear Regression:\n"
    "   - Mean Absolute Error (MAE): 50,413.43\n"
    "   - Mean Squared Error (MSE): 4.80e+09\n"
    "   - Root Mean Squared Error (RMSE): 69,297.72\n"
    "   - R-squared (R²): 0.6488\n\n"
    "2. Random Forest:\n"
    "   - Mean Absolute Error (MAE): 31,641.57\n"
    "   - Mean Squared Error (MSE): 2.38e+09\n"
    "   - Root Mean Squared Error (RMSE): 48,738.76\n"
    "   - R-squared (R²): 0.8263\n\n"
    "3. Neural Network:\n"
    "   - Mean Absolute Error (MAE): 43,290.98\n"
    "   - Mean Squared Error (MSE): 3.88e+09\n"
    "   - Root Mean Squared Error (RMSE): 62,308.14\n"
    "   - R-squared (R²): 0.7161\n\n"
    "4. XGBoost:\n"
    "   - Mean Absolute Error (MAE): 30,660.61\n"
    "   - Mean Squared Error (MSE): 2.20e+09\n"
    "   - Root Mean Squared Error (RMSE): 46,951.28\n"
    "   - R-squared (R²): 0.8388"
)
summary_content.font.size = Pt(18)
summary_content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add the Model Comparison R² chart
img_path = "plots/R2.png"  # Update the file path if necessary
slide.shapes.add_picture(img_path, Inches(0.5), Inches(4.5), Inches(8), Inches(3))

# Add a new slide for the Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Conclusion"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
title.alignment = PP_ALIGN.LEFT

# Add content for the conclusion
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Add paragraphs for unsupervised and supervised conclusions
unsupervised_paragraph = content_frame.add_paragraph()
unsupervised_paragraph.text = (
    "• **Unsupervised Analysis:**\n"
    "   - PCA effectively reduced dimensionality, capturing 83.6% of the variance with 3 components.\n"
    "   - KMeans clustering uncovered 3 distinct groups of housing data with moderate silhouette scores (~0.36).\n"
    "   - Insights from clustering revealed regional and economic patterns, which can aid targeted decision-making."
)
unsupervised_paragraph.font.size = Pt(18)
unsupervised_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

supervised_paragraph = content_frame.add_paragraph()
supervised_paragraph.text = (
    "\n• **Supervised Analysis:**\n"
    "   - XGBoost emerged as the best-performing model with R² = 0.8388 and the lowest MAE = 30,661.\n"
    "   - Random Forest also demonstrated strong performance (R² = 0.8263), emphasizing feature importance.\n"
    "   - Linear Regression, while interpretable, showed lower performance due to model simplicity.\n"
    "   - Neural Network provided reasonable performance but required careful tuning to avoid overfitting.\n"
    "   - Key features like 'median_income' and 'ocean_proximity' were consistently important across models."
)
supervised_paragraph.font.size = Pt(18)
supervised_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add a closing note
closing_paragraph = content_frame.add_paragraph()
closing_paragraph.text = (
    "\n• **Closing Note:**\n"
    "   - Unsupervised clustering offers valuable exploratory insights, while supervised models enable predictive precision.\n"
    "   - Combining both approaches provides a holistic view of the dataset and supports better decision-making."
)
closing_paragraph.font.size = Pt(18)
closing_paragraph.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue for emphasis
closing_paragraph.font.bold = True



# Save the presentation
prs.save("Housing_Price_Analysis_Presentation.pptx")

print("Presentation created successfully!")
